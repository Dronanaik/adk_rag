from pathlib import Path

import fitz
from docx import Document


def extract_from_pdf(file_path: str) -> str:
    document = fitz.open(file_path)

    pages = []

    for page_number, page in enumerate(document, start=1):
        text = page.get_text()

        if text.strip():
            pages.append(
                f"[Page {page_number}]\n{text}"
            )

    return "\n\n".join(pages)


def extract_from_docx(file_path: str) -> str:
    document = Document(file_path)

    paragraphs = []

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            paragraphs.append(paragraph.text)

    return "\n".join(paragraphs)


def extract_from_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read()


def extract_from_csv(file_path: str) -> str:
    return extract_from_text(file_path)


def extract_from_xlsx(file_path: str) -> str:
    """
    Extract text from an Excel (.xlsx / .xls) file using openpyxl.
    Each row is rendered as tab-separated values; sheets are separated by headings.
    """
    try:
        import openpyxl
    except ImportError:
        raise ValueError(
            "openpyxl is required to parse Excel files. "
            "Install it with: pip install openpyxl"
        )

    workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)

    sheet_texts = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        rows = []

        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            # Skip entirely empty rows
            if any(v.strip() for v in row_values):
                rows.append("\t".join(row_values))

        if rows:
            sheet_texts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

    workbook.close()
    return "\n\n".join(sheet_texts)


def extract_from_pptx(file_path: str) -> str:
    """
    Extract text from a PowerPoint (.pptx) file using python-pptx.
    Each slide's text frames are extracted in order.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError(
            "python-pptx is required to parse PowerPoint files. "
            "Install it with: pip install python-pptx"
        )

    presentation = Presentation(file_path)

    slides = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_lines = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                line = paragraph.text.strip()
                if line:
                    slide_lines.append(line)

        if slide_lines:
            slides.append(
                f"[Slide {slide_number}]\n" + "\n".join(slide_lines)
            )

    return "\n\n".join(slides)


def extract_text(file_path: str) -> str:
    """
    Extract text based on file extension.

    Supported formats:
      .pdf        — PyMuPDF (text-based PDFs only)
      .docx       — python-docx
      .xlsx/.xls  — openpyxl
      .pptx       — python-pptx
      .txt/.md/.csv/.json/.xml/.html/.htm/.log — plain text reader
    """

    extension = Path(file_path).suffix.lower()

    if extension == ".pdf":
        return extract_from_pdf(file_path)

    if extension == ".docx":
        return extract_from_docx(file_path)

    if extension in {".xlsx", ".xls"}:
        return extract_from_xlsx(file_path)

    if extension == ".pptx":
        return extract_from_pptx(file_path)

    if extension in {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".log",
    }:
        return extract_from_text(file_path)

    raise ValueError(
        f"Unsupported file type: '{extension}'. "
        "Supported types: .pdf, .docx, .xlsx, .xls, .pptx, "
        ".txt, .md, .csv, .json, .xml, .html, .htm, .log"
    )